"""Populate the OFFICIAL ISRO BAH 2026 idea-submission PPTX template with the
VESPER / PS7 content.

The organisers require the official template be used as-is. This script opens that
template and only *fills* it — it never alters the background art, theme, colours
or slide dimensions. Content slides are reused in place; two evidence slides
(Proof-of-Concept, Validation) are added as new slides that carry the template's
own content background so the look is identical.

Run:  .venv/bin/python hackathon/deck/build_pptx.py
Out:  hackathon/deck/BAH2026_PS7_idea_deck.pptx
"""
from __future__ import annotations
import os
from io import BytesIO

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(__file__)
TEMPLATE = os.path.join(HERE, "[Pub] ISRO BAH 2026 _ Idea Submission Template (2).pptx")
OUT = os.path.join(HERE, "BAH2026_PS7_idea_deck.pptx")
FIGS = os.path.join(HERE, "figs")
PFIGS = os.path.join(HERE, "..", "prototype", "figs")

FONT = "Arial"
NAVY = RGBColor(0x13, 0x29, 0x4B)
INK = RGBColor(0x1F, 0x27, 0x2E)
BODY = RGBColor(0x35, 0x3C, 0x43)
ACCENT = RGBColor(0x2C, 0x6F, 0xBB)
ORANGE = RGBColor(0xD8, 0x5A, 0x1E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LABEL = RGBColor(0xC9, 0xD6, 0xE5)
GREEN = RGBColor(0x1E, 0x84, 0x49)

# content-zone geometry (inches); slide is 10.0 x 5.625
TITLE_T, RULE_T, BODY_T, BODY_B = 0.60, 1.24, 1.44, 5.42


# ---------------------------------------------------------------- helpers ----
def _run(p, text, size, color, bold=False, italic=False, font=FONT):
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color
    return r


def instruction_box(slide):
    """The single template text box on a content slide (not the background pic)."""
    for sh in slide.shapes:
        if sh.has_text_frame:
            return sh
    return None


def set_title(slide, text, shape=None):
    """Turn the template instruction box (or a new box) into a styled slide title."""
    if shape is None:
        shape = slide.shapes.add_textbox(Inches(0.4), Inches(TITLE_T), Inches(9.2), Inches(0.6))
    shape.left, shape.top = Inches(0.4), Inches(TITLE_T)
    shape.width, shape.height = Inches(9.2), Inches(0.62)
    tf = shape.text_frame
    tf.word_wrap = True
    try:
        tf.auto_size = None
    except Exception:
        pass
    tf.clear()
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    size = 25.0
    max_w = 8.8  # keep the title on one line inside the 9.2in box
    if len(text) * 0.0072 * size > max_w:
        size = max(16.0, max_w / (len(text) * 0.0072))
    _run(p, text, size, NAVY, bold=True)
    # accent rule under the title
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(RULE_T), Inches(9.16), Inches(0.028))
    rule.fill.solid(); rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()
    rule.shadow.inherit = False
    return shape


def add_bullets(slide, items, left, top, width, height, size=13, gap=8, lead_color=NAVY):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap); p.line_spacing = 1.05
        _run(p, "•  ", size, ACCENT, bold=True)
        if isinstance(item, tuple):
            lead, rest = item
            _run(p, lead, size, lead_color, bold=True)
            if rest:
                _run(p, " — " + rest, size, BODY)
        else:
            _run(p, item, size, BODY)
    return tb


def add_caption(slide, lines, left, top, width, size=11, align=PP_ALIGN.CENTER, gap=5):
    box_h = max(0.4, min(1.2, 5.52 - top))  # never cross the slide bottom / footer
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(box_h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(gap); p.line_spacing = 1.05
        text, color, bold = ln if isinstance(ln, tuple) else (ln, BODY, False)
        _run(p, text, size, color, bold=bold)
    return tb


def add_fig(slide, path, left, top, maxw, maxh):
    """Place image scaled to fit (aspect-preserving) inside the box, centred."""
    iw, ih = Image.open(path).size
    ar = iw / ih
    w = maxw; h = w / ar
    if h > maxh:
        h = maxh; w = h * ar
    x = left + (maxw - w) / 2
    y = top + (maxh - h) / 2
    slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    return x, y, w, h


# ----------------------------------------------------------------- slides ----
def build_title(slide):
    fields = [("Team Name", "VESPER", 16),
              ("Team Leader Name", "Ansul Suryawanshi", 16),
              ("Problem Statement", "PS7 — AI-enabled Detection of Exoplanets from Noisy Astronomical Light Curves", 11)]
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        for key, val, sz in fields:
            if t.startswith(key):
                tf = sh.text_frame; tf.word_wrap = True
                try:
                    tf.auto_size = None
                except Exception:
                    pass
                tf.clear()
                p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
                _run(p, f"{key}:  ", sz, RGBColor(0x5A, 0x61, 0x68))   # band is white → dark label
                _run(p, val, sz + 0.5, NAVY, bold=True)                 # dark value
                break
    tb = slide.shapes.add_textbox(Inches(0.34), Inches(5.16), Inches(9.3), Inches(0.4))
    tf = tb.text_frame; tf.word_wrap = True
    _run(tf.paragraphs[0],
         "VESPER — Validation Engine for Stellar Photometric Evidence & Recovery",
         10.5, ACCENT, italic=True)


def build_team(slide):
    heading = table_shape = None
    for sh in slide.shapes:
        if sh.has_table:
            table_shape = sh
        elif sh.has_text_frame and sh.text_frame.text.strip().startswith("Team Members"):
            heading = sh
    heading.left, heading.top = Inches(0.85), Inches(0.72)
    heading.width, heading.height = Inches(8.3), Inches(0.55)
    tf = heading.text_frame; tf.word_wrap = True
    try:
        tf.auto_size = None
    except Exception:
        pass
    tf.clear()
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, "Team Members", 24, NAVY, bold=True)

    table_shape.left, table_shape.top, table_shape.width = Inches(0.85), Inches(1.55), Inches(8.3)
    tbl = table_shape.table
    tbl.columns[0].width = Inches(4.15); tbl.columns[1].width = Inches(4.15)
    tbl.rows[0].height = Inches(1.55); tbl.rows[1].height = Inches(1.55)
    members = [
        ("Team Leader", "Ansul Suryawanshi", "Indira Gandhi National Open University (IGNOU)"),
        ("Team Member-1", "Riddhi Jain", "Indira Gandhi National Open University (IGNOU)"),
        ("Team Member-2", "Samiksha Choudhary", "Priyadarshini College of Engineering (Hingna), Nagpur"),
        ("Team Member-3", "— optional (BAH allows a team of 3–4)", ""),
    ]
    cells = [tbl.cell(0, 0), tbl.cell(0, 1), tbl.cell(1, 0), tbl.cell(1, 1)]
    for cell, (role, name, college) in zip(cells, members):
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_top = Pt(4); cell.margin_bottom = Pt(4)
        tf = cell.text_frame; tf.word_wrap = True; tf.clear()
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        _run(p, role, 14, ACCENT, bold=True)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(3)
        _run(p2, name, 12.5, INK, bold=True)
        if college:
            p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
            _run(p3, college, 10.5, BODY)


def build_opportunity(slide):
    set_title(slide, "Opportunity", instruction_box(slide))
    add_bullets(slide, [
        ("Different from existing ideas",
         "most pipelines either detect OR classify, often with uncalibrated scores. TLS finds a periodic dip; VESPER "
         "goes further — it says WHAT the dip is (transit / eclipse / blend / other) with calibrated confidence, on a "
         "spine whose recall is non-inferior to full TLS."),
        ("How it solves the problem",
         "detect periodic dips → classify with physics + shape discriminators → attach SNR + calibrated confidence "
         "→ fit period / depth / duration → visualise."),
        ("USP",
         "hybrid physics + deep learning (interpretable, small-data robust); calibrated confidence (conformal); "
         "physics decides detection (recall-first); light-curve blend tells now, pixel-level centroid / difference-imaging "
         "in Round-2; extends a validated, pre-registered spine."),
        ("Already demonstrated on fresh MAST data",
         "Pi Mensae c blindly recovered (P = 6.262 d, 289 ppm, SDE 12.3); TIC 100029948 flagged as an eclipsing binary "
         "(depth ~25%, odd–even 0.23)."),
    ], 0.5, BODY_T, 9.0, 3.8, size=13, gap=9)


def build_features(slide):
    set_title(slide, "Features", instruction_box(slide))
    add_bullets(slide, [
        ("Robust periodic-dip detection", "in noisy, crowded fields"),
        ("4-class AI classification", "transit / eclipse / blend / other (physics features now; + CNN in Round-2)"),
        ("SNR / significance", "for every detected event"),
        ("Transit parameter fit", "period, depth, duration (+ uncertainties)"),
        ("Calibrated confidence", "+ annotated visualisations"),
        ("Reproducible", "recall non-inferior to full TLS"),
    ], 0.5, 1.62, 4.55, 3.7, size=13, gap=10)
    add_fig(slide, os.path.join(FIGS, "four_class_concept.png"), 5.05, 1.5, 4.6, 3.75)


def build_flow(slide):
    set_title(slide, "Process Flow — mirrors the PS7 five-step methodology", instruction_box(slide))
    add_fig(slide, os.path.join(FIGS, "process_flow.png"), 0.35, 1.55, 9.3, 3.6)


def build_char(slide):
    set_title(slide, "Characterization — trapezoid shape fit (PS7 step 03)", instruction_box(slide))
    _, y, _, h = add_fig(slide, os.path.join(PFIGS, "characterization.png"), 0.4, 1.40, 9.2, 3.02)
    add_caption(slide, [
        ("Depth is nearly equal for the EB and the planet — the trapezoid SHAPE separates them:", INK, False),
        ("EB flat-fraction 0.26 (V-shape → false positive)   vs   planet flat-fraction 0.67 (U-shape) — on real MAST data.",
         ACCENT, True),
        ("Depth → planet radius (planet vs massive planet);   shape U vs V → planet vs eclipsing binary.", INK, False),
    ], 0.4, y + h + 0.06, 9.2, size=10.5)


def build_arch(slide):
    set_title(slide, "Architecture", instruction_box(slide))
    add_fig(slide, os.path.join(FIGS, "arch_diagram.png"), 0.4, 1.5, 9.2, 3.75)


def build_tech(slide):
    set_title(slide, "Technologies", instruction_box(slide))
    add_bullets(slide, [
        ("Core scientific Python", "numpy · scipy · pandas · matplotlib"),
        ("Astronomy", "lightkurve · astropy · astroquery / MAST · wotan · transitleastsquares · batman"),
        ("Machine learning", "scikit-learn HistGBT (physics branch, built) · PyTorch CNN + conformal / MAPIE (Round-2)"),
        ("Method", "injection–recovery into real LCs · leakage-safe group CV · bootstrap CIs · recall parity vs full TLS"),
        ("All open-source", "no specialised or licensed software; data is public (MAST)"),
    ], 0.5, BODY_T, 9.0, 3.8, size=13.5, gap=11)


def build_cost(slide):
    set_title(slide, "Estimated Implementation Cost", instruction_box(slide))
    add_bullets(slide, [
        ("Software: ₹0", "entirely open-source stack"),
        ("Data: ₹0", "public TESS archive (MAST); one sector ~40–55 GB raw, kept slim by condition-then-discard"),
        ("Compute (TESS)", "standard laptops for development + a GPU session for CNN training; optional free cloud credits"),
        ("Scales to Kepler / K2 (est.)", "~200k long-cadence stars; a full transit search ≈ 5,000–10,000 CPU-core-hours "
         "— a short cloud burst on MAST / AWS Open Data (no egress). Data still ₹0; light curves ~few-hundred GB."),
        ("Net", "TESS negligible, Kepler a modest cloud burst — feasibility is high; the detection spine already runs on real data"),
    ], 0.5, BODY_T, 9.0, 3.95, size=12.5, gap=9)


def build_poc(slide):
    set_title(slide, "Proof of Concept — fresh MAST data (validated 2026-06-26)")
    _, y, _, h = add_fig(slide, os.path.join(PFIGS, "eb_vs_planet.png"), 0.35, 1.44, 9.3, 2.98)
    add_caption(slide, [
        ("Left: TIC 100029948 — eclipsing binary flagged by physics features (depth ~25%, V-shape, odd–even 0.23).", INK, False),
        ("Right: Pi Mensae c — blindly recovered (P = 6.262 d vs lit. 6.268; depth 289 ppm vs lit. ~315; SDE 12.3).", INK, False),
        ("The detection spine works end-to-end today; the physics-branch classifier is validated next (CNN + curated labels in Round-2).", ACCENT, True),
    ], 0.4, y + h + 0.05, 9.2, size=11)


def build_ablation(slide):
    set_title(slide, "Feature ablation — no single feature suffices (synthetic labels)")
    _, y, _, h = add_fig(slide, os.path.join(PFIGS, "ablation.png"), 1.05, 1.42, 7.9, 3.05)
    add_caption(slide, [
        ("Leakage-safe 5-fold group CV (no injection host shared train/test): depth-only 0.62, shape-only 0.53 — no "
         "single family exceeds 0.72; the full physics set reaches macro-F1 0.83.", INK, False),
        ("Quantifies the committee's point — depth does not discriminate; the combined physics features do.", ACCENT, True),
    ], 0.4, y + h + 0.05, 9.2, size=10.5)


def build_classifier(slide):
    set_title(slide, "Classifier — leakage-safe proof of path (synthetic labels)")
    _, y, _, h = add_fig(slide, os.path.join(PFIGS, "classifier_eval.png"), 0.3, 1.42, 9.4, 3.02)
    add_caption(slide, [
        ("Group-CV out-of-fold (n = 480): accuracy 0.83, macro-F1 0.83 (95% CI 0.80–0.86). Eclipses & systematics are "
         "separated cleanly; the residual is shallow transit vs blend.", INK, False),
        ("Honest scope: synthetic labels validate the pipeline, not real-world accuracy — Round-2 uses the organisers' "
         "curated set + pixel-level features.", ACCENT, True),
    ], 0.4, y + h + 0.05, 9.2, size=10.5)


def build_failure(slide):
    set_title(slide, "The hard case & Round-2 roadmap (transit vs blend)")
    _, y, _, h = add_fig(slide, os.path.join(PFIGS, "failure_case.png"), 0.95, 1.42, 8.1, 3.05)
    add_caption(slide, [
        ("At planet depth, on the same real host + noise, a planet (U) and a blended EB (V + faint secondary/odd-even) "
         "are nearly identical from the light curve alone.", INK, False),
        ("This is the residual confusion — and precisely why Round-2 adds pixel-level centroid / difference-imaging features.",
         ACCENT, True),
    ], 0.4, y + h + 0.05, 9.2, size=10.5)


def add_content_slide(prs, layout, bg_blob):
    slide = prs.slides.add_slide(layout)
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)  # drop inherited slide-number placeholder
    slide.shapes.add_picture(BytesIO(bg_blob), 0, 0, width=prs.slide_width, height=prs.slide_height)
    return slide


def main():
    prs = Presentation(TEMPLATE)
    slides = list(prs.slides)

    build_title(slides[0])
    build_team(slides[1])
    build_opportunity(slides[2])
    build_features(slides[3])
    build_flow(slides[4])
    build_char(slides[5])          # template "Wireframes/Mock (optional)" slot
    build_arch(slides[6])
    build_tech(slides[7])
    build_cost(slides[8])
    # slides[9] = closing "Thank You" (left untouched)

    layout = slides[2].slide_layout          # the BLANK content layout
    bg_blob = slides[2].shapes[0].image.blob  # image1 (white content background)
    closing_id = list(prs.slides._sldIdLst)[9]

    build_poc(add_content_slide(prs, layout, bg_blob))
    build_ablation(add_content_slide(prs, layout, bg_blob))
    build_classifier(add_content_slide(prs, layout, bg_blob))
    build_failure(add_content_slide(prs, layout, bg_blob))

    # keep the "Thank You" slide last
    prs.slides._sldIdLst.remove(closing_id)
    prs.slides._sldIdLst.append(closing_id)

    prs.save(OUT)
    print(f"deck -> {OUT}  ({os.path.getsize(OUT)/1e6:.2f} MB, {len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
